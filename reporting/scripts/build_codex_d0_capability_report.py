#!/usr/bin/env python3
"""Build the Codex D0 analysis-reproduction capability report slide deck."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/ubuntu/agent-network")
OUT = ROOT / "repos/general-codex-output/support/slides/codex_d0_capability_report.pptx"
FIG = ROOT / "repos/general-codex-output/support/figures"

NAVY = RGBColor(20, 34, 50)
INK = RGBColor(32, 43, 54)
MUTED = RGBColor(91, 103, 112)
BLUE = RGBColor(37, 99, 163)
TEAL = RGBColor(26, 132, 128)
ORANGE = RGBColor(210, 111, 42)
RED = RGBColor(183, 62, 62)
GREEN = RGBColor(76, 145, 92)
PURPLE = RGBColor(112, 86, 165)
LIGHT = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(215, 222, 230)


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def slide(prs: Presentation, bg=LIGHT):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    return s


def textbox(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def header(s, title, kicker=None, num=None, dark=False):
    color = WHITE if dark else NAVY
    sub = RGBColor(210, 220, 230) if dark else ORANGE
    if kicker:
        textbox(s, Inches(0.65), Inches(0.28), Inches(9.8), Inches(0.22), kicker.upper(), 8, sub, True)
    textbox(s, Inches(0.65), Inches(0.52), Inches(10.9), Inches(0.5), title, 22, color, True)
    if num is not None:
        textbox(s, Inches(12.1), Inches(0.42), Inches(0.5), Inches(0.3), str(num), 8, sub, True, PP_ALIGN.RIGHT)
    if not dark:
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.1), Inches(12.0), Inches(0.02))
        rule.fill.solid()
        rule.fill.fore_color.rgb = LINE
        rule.line.fill.background()


def footer(s):
    textbox(
        s,
        Inches(0.65),
        Inches(7.08),
        Inches(11.0),
        Inches(0.2),
        "Codex-assisted D0 AN reproduction stress test | prepared for discussion with Gian",
        7,
        MUTED,
    )


def bullets(s, items, x, y, w, h, size=14, color=INK, spacing=5):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.level = 0
    return box


def card(s, x, y, w, h, title, items, accent=BLUE, size=12):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = WHITE
    r.line.color.rgb = LINE
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    textbox(s, x + Inches(0.22), y + Inches(0.16), w - Inches(0.35), Inches(0.34), title, 13, accent, True)
    bullets(s, items, x + Inches(0.22), y + Inches(0.62), w - Inches(0.38), h - Inches(0.75), size)


def badge(s, x, y, text, color):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.85), Inches(0.35))
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    textbox(s, x + Inches(0.08), y + Inches(0.07), Inches(1.68), Inches(0.18), text, 8, WHITE, True, PP_ALIGN.CENTER)


def image_fit(s, path, x, y, w, h, caption=None):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    ww, hh = int(iw * scale), int(ih * scale)
    pic = s.shapes.add_picture(str(path), x + (w - ww) / 2, y, width=ww, height=hh)
    if caption:
        textbox(s, x, y + hh + Inches(0.08), w, Inches(0.3), caption, 8, MUTED, False, PP_ALIGN.CENTER)
    return pic


def flow(s, labels, x, y, w, h, colors):
    gap = Inches(0.12)
    bw = (w - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        bx = x + i * (bw + gap)
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, bw, h)
        r.fill.solid()
        r.fill.fore_color.rgb = colors[i % len(colors)]
        r.line.fill.background()
        textbox(s, bx + Inches(0.08), y + Inches(0.2), bw - Inches(0.16), h - Inches(0.25), label, 11, WHITE, True, PP_ALIGN.CENTER)


def bar_chart(s, data, x, y, w, h, colors, max_value=None, title=None):
    if title:
        textbox(s, x, y - Inches(0.35), w, Inches(0.25), title, 11, NAVY, True)
    maxv = max_value or max(v for _, v in data)
    label_w = Inches(2.35)
    bar_w = w - label_w - Inches(0.4)
    row_h = h / len(data)
    for i, (label, value) in enumerate(data):
        yy = y + i * row_h
        textbox(s, x, yy + Inches(0.04), label_w, Inches(0.22), label, 9, INK, False)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + label_w, yy + Inches(0.05), bar_w, Inches(0.18))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(228, 234, 241)
        bg.line.fill.background()
        fill = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + label_w, yy + Inches(0.05), bar_w * (value / maxv), Inches(0.18))
        fill.fill.solid()
        fill.fill.fore_color.rgb = colors[i % len(colors)]
        fill.line.fill.background()
        textbox(s, x + label_w + bar_w + Inches(0.08), yy + Inches(0.035), Inches(0.35), Inches(0.2), str(value), 9, INK, True)


def title_slide(prs):
    s = slide(prs, NAVY)
    textbox(s, Inches(0.75), Inches(0.8), Inches(11.5), Inches(0.55), "Codex-Assisted D0 Analysis Reproduction", 30, WHITE, True)
    textbox(s, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.45), "Capabilities, Limits, and Workflow Lessons", 22, RGBColor(220, 230, 240), False)
    flow(
        s,
        ["Stress test", "Reproduction", "Automation", "Physics checks", "Limits"],
        Inches(0.8),
        Inches(3.0),
        Inches(11.7),
        Inches(0.95),
        [BLUE, TEAL, ORANGE, PURPLE, RED],
    )
    textbox(
        s,
        Inches(0.85),
        Inches(5.15),
        Inches(10.8),
        Inches(0.75),
        "This report is about where agentic analysis works and fails, not a final D0 physics result.",
        20,
        WHITE,
        True,
    )
    textbox(s, Inches(0.85), Inches(6.25), Inches(10.8), Inches(0.3), "Prepared as Evan Zhang's report about the Codex stress test", 10, RGBColor(200, 210, 220))


def build():
    prs = prs_new()
    title_slide(prs)
    n = 2

    s = slide(prs)
    header(s, "Bottom line", "Executive summary", n)
    card(s, Inches(0.75), Inches(1.45), Inches(3.8), Inches(4.55), "What worked", [
        "Codex built a reproducible D0 workflow from input data to diagnostic mass peaks.",
        "It wrote reusable scripts, ran local/CRAB-style workflows, produced plots, tables, and documentation.",
        "It could compare outputs to AN references and keep artifacts organized."
    ], GREEN)
    card(s, Inches(4.82), Inches(1.45), Inches(3.8), Inches(4.55), "What partially worked", [
        "Several selections were rederived or independently stress-tested.",
        "ZDC/HF/topology studies reached AN-compatible behavior, but not all choices are fully independent.",
        "The D+ test showed rapid channel generalization, but only as a peak-finding diagnostic."
    ], ORANGE)
    card(s, Inches(8.9), Inches(1.45), Inches(3.8), Inches(4.55), "What did not fully work", [
        "The full AN was not reproduced end-to-end.",
        "Final corrections, systematics, trigger efficiency treatment, and final yield model remain incomplete or AN-guided.",
        "Codex still needs expert-supplied validation targets and physics judgment."
    ], RED)
    footer(s); n += 1

    s = slide(prs)
    header(s, "The stress test", "What Codex was actually asked to do", n)
    flow(s, ["MINIAOD", "forest", "event skim", "D0 skim", "mass peak", "AN-like report"], Inches(0.85), Inches(1.45), Inches(11.6), Inches(0.75), [BLUE, TEAL, ORANGE])
    bullets(s, [
        "Start from official 2023 PbPb UPC inputs where possible, not only final AN products.",
        "Recreate the analysis chain pragmatically: use AN/TWiki/professor input as validation and guidance where needed.",
        "Build enough infrastructure that future channels, such as D+, can reuse the same path.",
        "Document what is reproduced, what is proxy-only, and what is missing."
    ], Inches(1.05), Inches(2.75), Inches(11.1), Inches(2.0), 18)
    badge(s, Inches(1.05), Inches(5.45), "not a final paper result", RED)
    badge(s, Inches(3.05), Inches(5.45), "serious workflow test", BLUE)
    badge(s, Inches(5.05), Inches(5.45), "physics supervision required", ORANGE)
    footer(s); n += 1

    s = slide(prs)
    header(s, "What Codex handled well", "Automation and reproducibility", n)
    card(s, Inches(0.75), Inches(1.35), Inches(3.65), Inches(4.95), "Execution", [
        "Generated analysis scripts and wrappers.",
        "Ran local tests, CRAB-style productions, plotting, and report generation.",
        "Repeated workflows without relying on manual copy-paste."
    ], BLUE)
    card(s, Inches(4.85), Inches(1.35), Inches(3.65), Inches(4.95), "Organization", [
        "Kept logs, outputs, manifests, and plots durable.",
        "Promoted selected outputs into Overleaf/GitHub/Drive artifacts.",
        "Built reusable scripts instead of one-off terminal history."
    ], TEAL)
    card(s, Inches(8.95), Inches(1.35), Inches(3.65), Inches(4.95), "Validation", [
        "Compared results to AN-style references.",
        "Classified outputs as partial/proxy/missing rather than claiming success.",
        "Used data/MC/control samples to test cut choices."
    ], GREEN)
    footer(s); n += 1

    s = slide(prs)
    header(s, "Recreated chain reached a D0 mass peak", "Evidence: selected Kpi spectrum", n)
    image_fit(s, FIG / "an_proxy_plots/mass_fit_all_selected.png", Inches(0.75), Inches(1.35), Inches(6.4), Inches(4.9), "Inclusive selected D0 diagnostic mass fit")
    card(s, Inches(7.55), Inches(1.55), Inches(4.9), Inches(4.25), "What this demonstrates", [
        "The chain produced a visible D0 -> Kpi peak after event and candidate selections.",
        "The inclusive diagnostic fit found mean 1.862 GeV and sigma 12.7 MeV.",
        "This is raw-yield QA, not a corrected cross section."
    ], BLUE, 13)
    footer(s); n += 1

    s = slide(prs)
    header(s, "Binned mass fits exist, but they are not the final AN model", "Evidence: pT / rapidity bins", n)
    image_fit(s, FIG / "an_proxy_plots/mass_fit_grid_selected.png", Inches(0.65), Inches(1.25), Inches(7.0), Inches(4.4), "D0 mass fits in analysis-like bins")
    card(s, Inches(7.95), Inches(1.35), Inches(4.75), Inches(4.7), "Blunt status", [
        "6/6 diagnostic binned fits return OK status.",
        "One high-|y|, high-pT bin has low statistics and a >10 MeV mass shift warning.",
        "The AN final signal model, correction chain, and systematics are not fully rebuilt."
    ], ORANGE, 13)
    footer(s); n += 1

    s = slide(prs)
    header(s, "ZDC and HF selections were investigated with control data", "Cut derivation behavior", n)
    image_fit(s, FIG / "zdc_1n_thresholds/zdc_plus_1n_threshold.png", Inches(0.55), Inches(1.35), Inches(5.9), Inches(3.7), "ZDC 1n threshold scan")
    image_fit(s, FIG / "hf_gap_derivation/hf_plus_threshold_efficiencies.png", Inches(6.8), Inches(1.35), Inches(5.9), Inches(3.7), "HF gap threshold scan")
    bullets(s, [
        "ZDC: data/EmptyBX checks support AN-like 1n thresholds, but exact calibration/backport details remain delicate.",
        "HF: blind threshold scan lands near AN values; low-track/high-track controls make the rule physically interpretable.",
    ], Inches(0.9), Inches(5.55), Inches(11.4), Inches(0.85), 14)
    footer(s); n += 1

    s = slide(prs)
    header(s, "Topology cuts were the strongest cut-derivation stress test", "Signal retention vs sideband/background rejection", n)
    image_fit(s, FIG / "topological_floor_independent_sideband/floor_profile_grid.png", Inches(0.55), Inches(1.25), Inches(6.1), Inches(3.0), "Displaced-vertex floor scan")
    image_fit(s, FIG / "topological_pointing_independent_sideband/pointing_plateau_grid.png", Inches(6.85), Inches(1.25), Inches(5.9), Inches(3.0), "Pointing-angle plateau scan")
    card(s, Inches(0.9), Inches(4.8), Inches(11.25), Inches(1.35), "Interpretation", [
        "Codex could construct an a-priori protocol for sidebands, DsvpvSig/Dchi2cl floors, and pointing scans.",
        "Exact AN pointing rows are accepted by the plateau rule, but a fully AN-free tie-break is not yet frozen."
    ], PURPLE, 12)
    footer(s); n += 1

    s = slide(prs)
    header(s, "How independent are the selections?", "Blunt classification", n)
    bar_chart(s, [("independent", 2), ("partially independent", 10), ("AN specification", 3)], Inches(0.9), Inches(1.65), Inches(5.5), Inches(2.2), [GREEN, ORANGE, RED], 10, "Selection-definition audit")
    card(s, Inches(6.9), Inches(1.35), Inches(5.45), Inches(4.55), "What this means", [
        "Independent: sideband definition and topology floor derivation reached strongest status.",
        "Partially independent: most event/candidate selections have physics/control-sample justification but still match AN conventions.",
        "AN specification: physics target, trigger mapping, and final mass-fit model are not independently rebuilt."
    ], ORANGE, 13)
    footer(s); n += 1

    s = slide(prs)
    header(s, "AN coverage is broad but not complete", "Concrete reproduction status", n)
    bar_chart(s, [("partial targets", 66), ("missing targets", 26), ("external/reference", 3)], Inches(0.85), Inches(1.55), Inches(5.7), Inches(2.25), [ORANGE, RED, MUTED], 66, "Figure/table target plan: 95 targets")
    bar_chart(s, [("partial sections", 9), ("strong partial", 3), ("partial scaffold", 3)], Inches(0.85), Inches(4.3), Inches(5.7), Inches(1.9), [ORANGE, GREEN, BLUE], 9, "Section-level coverage: 15 sections")
    card(s, Inches(7.1), Inches(1.45), Inches(5.2), Inches(4.7), "No overclaim", [
        "No AN section should be called fully reproduced yet.",
        "Strong partials: MC samples, offline event selection, D-meson selection/signal extraction.",
        "Partial scaffolds: cross section, final results, theory comparison."
    ], RED, 13)
    footer(s); n += 1

    s = slide(prs)
    header(s, "What did not work yet", "Physics-facing limitations only", n)
    card(s, Inches(0.75), Inches(1.35), Inches(3.8), Inches(4.95), "Full AN reproduction", [
        "The report/plots cover many AN-like targets, but the full note is not regenerated with final-quality numbers.",
        "Several figures are proxy or scaffolded rather than final."
    ], RED)
    card(s, Inches(4.82), Inches(1.35), Inches(3.8), Inches(4.95), "Corrections and systematics", [
        "Trigger efficiencies, acceptance/reco efficiencies, prompt fractions, luminosity normalization, and systematic variations are incomplete.",
        "Raw-yield QA exists; final cross sections do not."
    ], RED)
    card(s, Inches(8.9), Inches(1.35), Inches(3.8), Inches(4.95), "MC/ZDC constraints", [
        "Official MC was useful for topology and signal checks, but not sufficient for ZDC threshold derivation.",
        "ZDC calibration and neutron-response details remain a sensitive expert-controlled piece."
    ], RED)
    footer(s); n += 1

    s = slide(prs)
    header(s, "What user and expert input was actually needed", "Not everything was autonomous", n)
    card(s, Inches(0.75), Inches(1.35), Inches(3.8), Inches(4.8), "Document pointers", [
        "AN PDF and scope of reproduction.",
        "TWiki dataset locations and official reference products.",
        "Clarification that skims/forests should be recreated, not assumed."
    ], BLUE)
    card(s, Inches(4.82), Inches(1.35), Inches(3.8), Inches(4.8), "Physics judgment", [
        "Professor guidance on ZDC, HF, and topology-cut logic.",
        "Which uncertainties matter and which control samples are legitimate.",
        "What counts as acceptable agreement with the AN."
    ], ORANGE)
    card(s, Inches(8.9), Inches(1.35), Inches(3.8), Inches(4.8), "Validation targets", [
        "AN cut values as posthoc reference points.",
        "Expected mass peaks and binning.",
        "Which missing pieces are true blockers versus lower-priority scaffolds."
    ], TEAL)
    footer(s); n += 1

    s = slide(prs)
    header(s, "Generalization test: D+ peak finding", "Same infrastructure, different channel", n)
    image_fit(s, FIG / "dplus_peak/dplus_mass_dplus_medium_prelim59_toposcan_20260623T1230Z_tight_topology.png", Inches(0.75), Inches(1.35), Inches(6.4), Inches(4.55), "D+ -> Kpipi diagnostic peak, 59-file checkpoint")
    card(s, Inches(7.55), Inches(1.45), Inches(4.9), Inches(4.45), "What this shows", [
        "Codex modified the reconstruction path from D0 -> Kpi to D+ -> Kpipi.",
        "The 59-file checkpoint gave a visible peak near the D+ mass.",
        "This is a diagnostic peak only: D+ cuts, efficiencies, and systematics are not final."
    ], GREEN, 13)
    footer(s); n += 1

    s = slide(prs)
    header(s, "How Codex should be used for this kind of analysis", "Best use cases", n)
    card(s, Inches(0.75), Inches(1.35), Inches(3.8), Inches(4.8), "Good fit", [
        "Build repeatable analysis plumbing.",
        "Generate cut scans and diagnostic plots.",
        "Maintain logs, manifests, and reports.",
        "Compare outputs against references."
    ], GREEN)
    card(s, Inches(4.82), Inches(1.35), Inches(3.8), Inches(4.8), "Use with oversight", [
        "Selection derivation.",
        "MC/data control-region interpretation.",
        "Model choice and systematic definitions.",
        "Claims of final physics agreement."
    ], ORANGE)
    card(s, Inches(8.9), Inches(1.35), Inches(3.8), Inches(4.8), "Bad fit alone", [
        "Final analysis approval.",
        "Unvalidated detector-calibration decisions.",
        "Inferring missing collaboration-specific conventions.",
        "Replacing expert review."
    ], RED)
    footer(s); n += 1

    s = slide(prs)
    header(s, "Recommended next iteration", "What would make the stress test stronger", n)
    bullets(s, [
        "Freeze a final input registry: data, MC, EmptyBX, trigger menus, luminosity, and official validation products.",
        "Build a final-yield fitter with predeclared signal/background alternatives and mass-shape bias checks.",
        "Turn correction inputs into reproducible producers: trigger efficiency, acceptance/reco efficiency, prompt fraction, EMD/ZDC effects.",
        "Run systematic variations as a single scripted workflow, not as isolated plots.",
        "Keep Codex in the loop for automation and cross-checks, but require expert signoff on every physics assumption."
    ], Inches(1.0), Inches(1.45), Inches(11.4), Inches(4.7), 17)
    footer(s); n += 1

    s = slide(prs, NAVY)
    header(s, "Final assessment", "Blunt conclusion", n, dark=True)
    textbox(
        s,
        Inches(0.9),
        Inches(1.55),
        Inches(11.3),
        Inches(1.2),
        "Codex is already useful as a persistent analysis operator: it can build workflows, run scans, generate plots, compare references, and document what happened.",
        22,
        WHITE,
        True,
    )
    textbox(
        s,
        Inches(0.9),
        Inches(3.2),
        Inches(11.3),
        Inches(1.25),
        "It is not yet a standalone physics analyst. The weak points are not typing code; they are final model choices, detector-specific calibration, correction chains, and deciding what agreement is scientifically sufficient.",
        21,
        RGBColor(225, 232, 240),
        True,
    )
    textbox(
        s,
        Inches(0.9),
        Inches(5.25),
        Inches(11.3),
        Inches(0.65),
        "Best use: let Codex do the repetitive, auditable work; keep human physicists responsible for the assumptions.",
        20,
        RGBColor(255, 218, 170),
        True,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
