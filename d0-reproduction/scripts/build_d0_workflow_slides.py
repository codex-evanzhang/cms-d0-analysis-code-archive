#!/usr/bin/env python3
"""Build a concise slide deck for the D0 AN recreation workflow."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(24, 38, 55)
BLUE = RGBColor(36, 99, 163)
TEAL = RGBColor(36, 140, 132)
ORANGE = RGBColor(207, 112, 48)
GRAY = RGBColor(91, 102, 112)
LIGHT = RGBColor(244, 247, 250)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, left, top, width, height, text, size=22, color=NAVY, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_header(slide, title, kicker=None):
    if kicker:
        add_textbox(slide, Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.25), kicker, 9, ORANGE, True)
    add_textbox(slide, Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.55), title, 24, NAVY, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.15), Inches(11.9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(215, 222, 230)
    line.line.fill.background()


def add_footer(slide, number):
    add_textbox(
        slide,
        Inches(0.7),
        Inches(7.1),
        Inches(8.5),
        Inches(0.22),
        "Source: Barebones D0 Analysis Workflow / general-codex-output",
        7,
        GRAY,
        False,
    )
    add_textbox(slide, Inches(12.0), Inches(7.1), Inches(0.6), Inches(0.22), str(number), 8, GRAY, False)


def add_bullets(slide, items, left, top, width, height, size=16, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def add_card(slide, left, top, width, height, title, bullets, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(210, 218, 226)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_textbox(slide, left + Inches(0.22), top + Inches(0.15), width - Inches(0.35), Inches(0.32), title, 14, accent, True)
    add_bullets(slide, bullets, left + Inches(0.22), top + Inches(0.55), width - Inches(0.35), height - Inches(0.65), 10, NAVY)


def add_flow(slide, labels, top=Inches(2.4)):
    x0 = Inches(0.7)
    box_w = Inches(1.72)
    gap = Inches(0.18)
    colors = [BLUE, TEAL, ORANGE, BLUE, TEAL, ORANGE]
    for i, label in enumerate(labels):
        x = x0 + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        if i < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w - Inches(0.03), top + Inches(0.34), gap + Inches(0.06), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(170, 178, 186)
            arrow.line.fill.background()


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(slide, Inches(0.85), Inches(1.05), Inches(11.3), Inches(0.45), "D0 AN Recreation Workflow", 32, WHITE, True)
    add_textbox(slide, Inches(0.9), Inches(1.7), Inches(10.6), Inches(0.6), "High-level steps needed to get from collision data to a D0 mass peak", 21, RGBColor(220, 230, 240), False)
    add_flow(slide, ["Data", "Event selection", "D0 candidates", "Background rejection", "Mass fit", "Validation"], Inches(3.25))
    add_textbox(slide, Inches(0.9), Inches(5.25), Inches(10.8), Inches(0.5), "The goal is to understand what must be done, not the exact commands used to do it.", 18, WHITE, True)
    add_textbox(slide, Inches(0.9), Inches(6.05), Inches(10.8), Inches(0.35), "Generated by Codex from the Barebones D0 workflow document.", 10, RGBColor(200, 210, 220), False)
    return slide


def build_deck(outpath: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    slide_no = 2

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "The analysis chain in one view")
    add_flow(slide, ["Certified PbPb data", "UPC 0nXn events", "Good Kpi pairs", "Topology-selected D0", "Kpi mass peak", "AN closure"], Inches(1.8))
    add_bullets(
        slide,
        [
            "Each step removes a different class of background before the final mass fit.",
            "The AN tells us what behavior the reproduced chain should match.",
            "The useful question at each step is: what physical or detector effect is this selection controlling?",
        ],
        Inches(1.1),
        Inches(4.0),
        Inches(11.1),
        Inches(1.5),
        17,
    )
    add_footer(slide, slide_no)
    slide_no += 1

    slides = [
        (
            "1. Start from the correct collision sample",
            "Define the dataset before making any physics selections.",
            [
                ("Collision data", ["2023 PbPb UPC data", "good certified runs/lumisections", "trigger paths that cover the D0 pT bins"]),
                ("Reference samples", ["signal MC for D0 efficiency", "EmptyBX data for noise checks", "official AN products only for validation"]),
                ("Purpose", ["make sure later differences are physics/selection differences, not sample mismatches"]),
            ],
        ),
        (
            "2. Reconstruct the objects needed for D0",
            "Turn the event record into analysis-level quantities.",
            [
                ("Event objects", ["primary vertex", "event-quality flags", "forward calorimeter activity"]),
                ("UPC handles", ["ZDC energy on both beam sides", "HF maximum energy on both sides", "trigger decisions"]),
                ("Charm handles", ["charged tracks", "Kpi invariant-mass candidates", "secondary-vertex quantities"]),
            ],
        ),
        (
            "3. Select UPC-like 0nXn events",
            "Remove ordinary hadronic PbPb backgrounds before looking at D0 candidates.",
            [
                ("Event quality", ["require a good primary vertex", "reject halo/noise-like events", "suppress high-multiplicity contamination"]),
                ("ZDC topology", ["require exactly one neutron-like ZDC side", "require the opposite side to be 0n-like", "this defines Xn0n or 0nXn"]),
                ("HF rapidity gap", ["on the 0n/photon-going side, require no large HF energy deposit", "this rejects hadronic or beam-gas activity"]),
            ],
        ),
        (
            "4. Build plausible D0 -> Kpi candidates",
            "Keep track pairs that could be real D0 decays.",
            [
                ("Track pair", ["combine opposite-charge tracks", "test both K/pi mass assignments", "require tracks to be well measured"]),
                ("D0 phase space", ["keep candidates in the target pT and rapidity range", "use a broad mass window for fitting and sidebands"]),
                ("Why this matters", ["track quality removes fake pairs", "the broad mass region avoids sculpting the peak"]),
            ],
        ),
        (
            "5. Suppress random Kpi combinations",
            "Use decay geometry to separate displaced D0 decays from background.",
            [
                ("Displacement", ["require a secondary vertex clearly separated from the primary vertex", "larger decay-length significance is more D0-like"]),
                ("Vertex quality", ["require the two tracks to form a credible common decay vertex", "secondary-vertex probability controls this"]),
                ("Pointing", ["D0 momentum should point back to the primary vertex", "daughter opening angle should be compatible with a D0 decay"]),
            ],
        ),
        (
            "6. Choose cuts without tuning on the peak",
            "Use control information to justify the main selections.",
            [
                ("ZDC cuts", ["derive neutron thresholds from data and EmptyBX noise", "check that fake 1n assignment is negligible"]),
                ("HF gap cuts", ["pick thresholds that keep UPC-like events quiet", "check leakage using high-activity control regions"]),
                ("Topology cuts", ["use MC for signal retention", "use mass sidebands for combinatorial background", "validate against BDT-style cross-checks"]),
            ],
        ),
        (
            "7. Extract the D0 signal",
            "Fit the Kpi invariant-mass spectrum after selections.",
            [
                ("Mass plots", ["first make an inclusive sanity-check peak", "then split by pT and rapidity bins"]),
                ("Fit components", ["D0 signal shape", "smooth combinatorial background", "wrong-mass and peaking backgrounds where relevant"]),
                ("Result", ["raw D0 yields in each bin", "fit quality checks", "comparison to AN trends"]),
            ],
        ),
        (
            "8. Convert the peak into an AN-level result",
            "The mass peak is necessary but not the full measurement.",
            [
                ("Corrections", ["trigger and reconstruction efficiencies", "acceptance and selection efficiency", "prompt/nonprompt treatment"]),
                ("Normalization", ["luminosity", "branching fraction", "bin widths and event topology definitions"]),
                ("Systematics", ["cut variations", "fit-model variations", "correction uncertainties"]),
            ],
        ),
    ]

    accents = [BLUE, TEAL, ORANGE]
    for title, subtitle, cards in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header(slide, title, subtitle)
        for i, (card_title, bullets) in enumerate(cards):
            add_card(slide, Inches(0.8 + i * 4.15), Inches(1.65), Inches(3.75), Inches(4.6), card_title, bullets, accents[i % len(accents)])
        add_footer(slide, slide_no)
        slide_no += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "High-level closure checks")
    add_bullets(
        slide,
        [
            "The selected event sample should look UPC-like: one neutron side, one quiet side, and a clean HF gap.",
            "The D0 candidate sample should show a stable Kpi peak without needing peak-sculpting cuts.",
            "Each nontrivial threshold should have a control-region or MC/sideband justification.",
            "Binned mass fits should reproduce the AN's qualitative trends before moving to corrected cross sections.",
            "Any mismatch should be labeled as statistics-limited, input-limited, method-limited, or a real bug.",
        ],
        Inches(1.0),
        Inches(1.65),
        Inches(11.5),
        Inches(4.8),
        18,
    )
    add_footer(slide, slide_no)

    outpath.parent.mkdir(parents=True, exist_ok=True)
    prs.save(outpath)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--out",
        default="repos/general-codex-output/support/slides/d0_an_recreation_workflow_steps.pptx",
        help="Output .pptx path.",
    )
    args = parser.parse_args()
    build_deck(Path(args.out))
    print(args.out)


if __name__ == "__main__":
    main()
